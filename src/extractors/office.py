"""
Office document extraction module (Word, Excel, PowerPoint)
"""

import logging
from typing import Dict, Any, Optional, List
from pathlib import Path

try:
    from docx import Document as WordDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    from odf import text, teletype
    from odf.opendocument import load
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False

from .base import BaseExtractor

logger = logging.getLogger(__name__)


class OfficeExtractor(BaseExtractor):
    """Extract text and metadata from Office documents"""

    def __init__(self, config: Dict[str, Any] = None):
        """Initialize office document extractor

        Args:
            config: Configuration dictionary
        """
        super().__init__(config)
        self.extract_comments = config.get('extract_comments', True)
        self.extract_metadata = config.get('extract_metadata', True)
        self.extract_headers_footers = config.get('extract_headers_footers', True)
        self.extract_tables = config.get('extract_tables', True)
        self.preserve_formatting = config.get('preserve_formatting', False)

    def extract(self, source: str, **kwargs) -> Dict[str, Any]:
        """Extract text and metadata from office document

        Args:
            source: Path to office document
            **kwargs: Additional extraction options

        Returns:
            Dictionary with extracted text and metadata
        """
        if not self.validate_source(source):
            return {
                'text': '',
                'metadata': {},
                'error': f'Invalid source: {source}'
            }

        file_ext = Path(source).suffix.lower()

        try:
            if file_ext in ['.docx', '.doc']:
                result = self._extract_word(source, **kwargs)
            elif file_ext in ['.xlsx', '.xls']:
                result = self._extract_excel(source, **kwargs)
            elif file_ext in ['.pptx', '.ppt']:
                result = self._extract_powerpoint(source, **kwargs)
            elif file_ext in ['.odt', '.ods', '.odp']:
                result = self._extract_opendocument(source, **kwargs)
            else:
                return {
                    'text': '',
                    'metadata': self.extract_metadata(source),
                    'error': f'Unsupported file format: {file_ext}'
                }

            self.update_stats(
                success=True,
                bytes_processed=Path(source).stat().st_size
            )

            return result

        except Exception as e:
            logger.error(f"Error extracting office document {source}: {e}")
            self.update_stats(success=False)
            return {
                'text': '',
                'metadata': self.extract_metadata(source),
                'error': str(e)
            }

    def _extract_word(self, source: str, **kwargs) -> Dict[str, Any]:
        """Extract from Word documents"""
        if not DOCX_AVAILABLE:
            return {
                'text': '',
                'metadata': {},
                'error': 'python-docx not installed'
            }

        text_parts = []
        metadata = self.extract_metadata(source)
        tables_data = []

        doc = WordDocument(source)

        # Extract core properties
        if self.extract_metadata:
            props = doc.core_properties
            metadata.update({
                'title': props.title,
                'author': props.author,
                'subject': props.subject,
                'keywords': props.keywords,
                'created': props.created.isoformat() if props.created else None,
                'modified': props.modified.isoformat() if props.modified else None,
                'last_modified_by': props.last_modified_by,
            })

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                if self.preserve_formatting:
                    # Preserve basic formatting info
                    text_parts.append(self._format_paragraph(para))
                else:
                    text_parts.append(para.text)

        # Extract tables
        if self.extract_tables:
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables_data.append(table_data)

        # Extract headers and footers
        if self.extract_headers_footers:
            for section in doc.sections:
                if section.header:
                    header_text = ' '.join([p.text for p in section.header.paragraphs])
                    if header_text.strip():
                        text_parts.append(f"[Header: {header_text}]")
                if section.footer:
                    footer_text = ' '.join([p.text for p in section.footer.paragraphs])
                    if footer_text.strip():
                        text_parts.append(f"[Footer: {footer_text}]")

        # Extract comments
        if self.extract_comments and hasattr(doc, 'comments'):
            for comment in doc.comments:
                text_parts.append(f"[Comment: {comment.text}]")

        return {
            'text': '\n'.join(text_parts),
            'metadata': metadata,
            'tables': tables_data if tables_data else None
        }

    def _extract_excel(self, source: str, **kwargs) -> Dict[str, Any]:
        """Extract from Excel documents"""
        text_parts = []
        metadata = self.extract_metadata(source)
        all_data = []

        file_ext = Path(source).suffix.lower()

        if file_ext == '.xlsx' and OPENPYXL_AVAILABLE:
            wb = openpyxl.load_workbook(source, data_only=True)

            # Extract metadata
            if self.extract_metadata and hasattr(wb, 'properties'):
                props = wb.properties
                metadata.update({
                    'title': props.title,
                    'creator': props.creator,
                    'created': props.created.isoformat() if props.created else None,
                    'modified': props.modified.isoformat() if props.modified else None,
                })

            # Extract data from each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_data = {
                    'sheet_name': sheet_name,
                    'data': []
                }

                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        if cell.value is not None:
                            row_data.append(str(cell.value))
                    if row_data:
                        sheet_data['data'].append(row_data)

                if sheet_data['data']:
                    all_data.append(sheet_data)
                    # Convert to text
                    text_parts.append(f"\n[Sheet: {sheet_name}]")
                    for row in sheet_data['data']:
                        text_parts.append('\t'.join(row))

            wb.close()

        elif file_ext == '.xls' and XLRD_AVAILABLE:
            wb = xlrd.open_workbook(source)

            # Extract data from each sheet
            for sheet_idx in range(wb.nsheets):
                sheet = wb.sheet_by_index(sheet_idx)
                sheet_data = {
                    'sheet_name': sheet.name,
                    'data': []
                }

                for row_idx in range(sheet.nrows):
                    row_data = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        row_data.append(str(cell.value))
                    sheet_data['data'].append(row_data)

                if sheet_data['data']:
                    all_data.append(sheet_data)
                    # Convert to text
                    text_parts.append(f"\n[Sheet: {sheet.name}]")
                    for row in sheet_data['data']:
                        text_parts.append('\t'.join(row))

        else:
            return {
                'text': '',
                'metadata': metadata,
                'error': 'Excel libraries not available for this file type'
            }

        return {
            'text': '\n'.join(text_parts),
            'metadata': metadata,
            'sheets': all_data
        }

    def _extract_powerpoint(self, source: str, **kwargs) -> Dict[str, Any]:
        """Extract from PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            return {
                'text': '',
                'metadata': {},
                'error': 'python-pptx not installed'
            }

        text_parts = []
        metadata = self.extract_metadata(source)
        slides_data = []

        prs = Presentation(source)

        # Extract metadata
        if self.extract_metadata:
            props = prs.core_properties
            metadata.update({
                'title': props.title,
                'author': props.author,
                'subject': props.subject,
                'created': props.created.isoformat() if props.created else None,
                'modified': props.modified.isoformat() if props.modified else None,
                'slide_count': len(prs.slides)
            })

        # Extract text from each slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    slide_text.append(shape.text)
                elif shape.has_table:
                    # Extract table data
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(' | '.join(row_text))

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_text.append(f"[Notes: {notes}]")

            if slide_text:
                slides_data.append({
                    'slide_number': slide_idx + 1,
                    'text': '\n'.join(slide_text)
                })
                text_parts.append(f"\n[Slide {slide_idx + 1}]")
                text_parts.extend(slide_text)

        return {
            'text': '\n'.join(text_parts),
            'metadata': metadata,
            'slides': slides_data
        }

    def _extract_opendocument(self, source: str, **kwargs) -> Dict[str, Any]:
        """Extract from OpenDocument formats"""
        if not ODF_AVAILABLE:
            return {
                'text': '',
                'metadata': {},
                'error': 'odfpy not installed'
            }

        text_parts = []
        metadata = self.extract_metadata(source)

        try:
            doc = load(source)

            # Extract text based on document type
            file_ext = Path(source).suffix.lower()

            if file_ext == '.odt':
                # Extract text from ODT
                for element in doc.getElementsByType(text.P):
                    paragraph_text = teletype.extractText(element)
                    if paragraph_text.strip():
                        text_parts.append(paragraph_text)

            elif file_ext == '.ods':
                # Extract from ODS spreadsheet
                from odf.table import Table, TableRow, TableCell
                for table in doc.getElementsByType(Table):
                    table_name = table.getAttribute('name')
                    text_parts.append(f"\n[Table: {table_name}]")
                    for row in table.getElementsByType(TableRow):
                        row_data = []
                        for cell in row.getElementsByType(TableCell):
                            cell_text = teletype.extractText(cell)
                            row_data.append(cell_text)
                        if row_data:
                            text_parts.append('\t'.join(row_data))

            elif file_ext == '.odp':
                # Extract from ODP presentation
                for element in doc.getElementsByType(text.P):
                    paragraph_text = teletype.extractText(element)
                    if paragraph_text.strip():
                        text_parts.append(paragraph_text)

        except Exception as e:
            logger.error(f"Error extracting OpenDocument {source}: {e}")
            return {
                'text': '',
                'metadata': metadata,
                'error': str(e)
            }

        return {
            'text': '\n'.join(text_parts),
            'metadata': metadata
        }

    def _format_paragraph(self, paragraph) -> str:
        """Format paragraph with basic styling info"""
        result = []
        for run in paragraph.runs:
            text = run.text
            if run.bold:
                text = f"**{text}**"
            if run.italic:
                text = f"*{text}*"
            result.append(text)
        return ''.join(result)